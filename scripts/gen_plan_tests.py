# -*- coding: utf-8 -*-
"""Génère PLAN_TESTS_LaBarakat.pptx — plan de tests de l'application."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

VERT   = RGBColor(0x16, 0x65, 0x34)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
ROUGE  = RGBColor(0xDC, 0x26, 0x26)
GRIS   = RGBColor(0x64, 0x74, 0x8B)
NOIR   = RGBColor(0x0F, 0x17, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide_titre(titre, sous_titre):
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = titre
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = VERT
    box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.2))
    tf = box2.text_frame
    for i, line in enumerate(sous_titre.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20); p.font.color.rgb = GRIS
    return s

def slide_test(num, titre, objectif, etapes, attendu, badge_color=TEAL):
    s = prs.slides.add_slide(BLANK)
    # Bandeau numéro + titre
    head = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = head.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = f"TEST {num}  "
    r1.font.size = Pt(30); r1.font.bold = True; r1.font.color.rgb = badge_color
    r2 = p.add_run(); r2.text = titre
    r2.font.size = Pt(30); r2.font.bold = True; r2.font.color.rgb = NOIR
    # Objectif
    obj = s.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.6))
    p = obj.text_frame.paragraphs[0]
    p.text = "🎯 Objectif : " + objectif
    p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = GRIS
    # Étapes
    et = s.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(7.3), Inches(5.1))
    tf = et.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]; p0.text = "ÉTAPES"; p0.font.size = Pt(14); p0.font.bold = True; p0.font.color.rgb = TEAL
    for i, e in enumerate(etapes, 1):
        p = tf.add_paragraph()
        p.text = f"{i}.  {e}"
        p.font.size = Pt(15); p.font.color.rgb = NOIR
        p.space_after = Pt(6)
    # Résultat attendu
    att = s.shapes.add_textbox(Inches(8.2), Inches(1.95), Inches(4.6), Inches(5.1))
    tf = att.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]; p0.text = "✅ RÉSULTAT ATTENDU"; p0.font.size = Pt(14); p0.font.bold = True; p0.font.color.rgb = VERT
    for a in attendu:
        p = tf.add_paragraph()
        p.text = "• " + a
        p.font.size = Pt(14); p.font.color.rgb = NOIR
        p.space_after = Pt(6)
    return s

def slide_liste(titre, items, couleur=TEAL, taille=16):
    s = prs.slides.add_slide(BLANK)
    head = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = head.text_frame.paragraphs[0]
    p.text = titre
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = couleur
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6))
    tf = body.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.font.size = Pt(taille); p.font.color.rgb = NOIR
        p.space_after = Pt(10)
    return s

# ══ SLIDE 1 — Titre ══
slide_titre("🐾 La Barakat — Plan de tests",
            "Nouvelle version : cliniquebarakat-app-v2.pages.dev\nValidation des fonctionnalités livrées · Août 2026")

# ══ SLIDE 2 — Prérequis ══
slide_liste("⚙️ Prérequis avant de commencer", [
    "1.  Exécuter dans Supabase → SQL Editor : supabase/commandes_echeance.sql (le seul script encore en attente)",
    "2.  Vérifier que les fiches Médicaments ont un PRIX D'ACHAT renseigné (sinon marges faussées à 100 %)",
    "3.  Avoir au moins un médicament de catégorie « Vaccin » en stock (ex : Vaccin Rage)",
    "4.  Disposer de 2 comptes de test : un ADMIN et un VÉTÉRINAIRE (ou caissier)",
    "5.  Utiliser la nouvelle URL : https://cliniquebarakat-app-v2.pages.dev",
    "",
    "💡 Astuce : ouvrez la console navigateur (F12) pendant les tests — toute erreur rouge est à signaler.",
], couleur=ORANGE)

# ══ TEST 1 — Connexion & rôles ══
slide_test(1, "Connexion & Sécurité des rôles (RLS)",
    "Vérifier que chaque rôle ne voit que ce qui le concerne",
    ["Se connecter avec le compte ADMIN → vérifier que toutes les données chargent (patients, ventes, stock)",
     "Se déconnecter, se connecter avec le compte VÉTÉRINAIRE",
     "Vérifier l'accès aux pages : Patients, Consultations, Chirurgies, Agenda",
     "Tenter d'accéder aux pages Caisse, Dépenses, Fournisseurs, Comptes"],
    ["L'admin voit tout",
     "Le vétérinaire voit ses pages avec leurs données",
     "Les pages hors rôle affichent « Interdit » ou sont absentes du menu",
     "Aucune page blanche, aucune erreur console"])

# ══ TEST 2 — Caisse ══
slide_test(2, "Caisse — vente au comptoir",
    "Valider le calcul TTC, le stock par conditionnement et le reçu",
    ["Ouvrir la Caisse → noter le stock actuel d'un médicament en comprimés",
     "Vendre 2 « Plaquette (10 cp) » de ce médicament, encaisser en Espèces",
     "Imprimer le reçu et vérifier les montants (HT, TVA si active, TTC)",
     "Aller dans Médicaments → vérifier le stock"],
    ["Le total = 2 × prix plaquette",
     "Le reçu affiche le bon TOTAL TTC",
     "Le stock a baissé de 20 comprimés (2 × 10), pas de 2",
     "La vente apparaît dans l'historique avec le caissier"])

# ══ TEST 3 — Créances ══
slide_test(3, "Créances — paiement partiel",
    "Valider le suivi du restant dû et les versements",
    ["En Caisse, faire une vente de 5 000 F en mode « À crédit » au nom d'un client",
     "Ouvrir Créances → retrouver le client",
     "Cliquer « 💵 Versement » → saisir 2 000 F",
     "Vérifier le restant dû, puis verser les 3 000 F restants"],
    ["Après 2 000 F : statut « Partiellement payé », restant 3 000 F, mention « déjà payé : 2 000 F »",
     "Après le solde : statut « Payé », la créance disparaît",
     "Le stock est décrémenté au paiement final",
     "Le CA encaissé (Rapports) reflète les montants payés"], badge_color=ORANGE)

# ══ TEST 4 — Consultation avec traitements ══
slide_test(4, "Consultation — traitements facturés",
    "Le scénario « 1 ml = 100 F, 5 ml = 500 F »",
    ["Consultations → Nouvelle consultation → choisir un PATIENT via la suggestion",
     "Remplir le diagnostic (SOAP-A), mettre 3 000 F d'actes",
     "« 💉 Traitements » → + Ajouter → choisir l'antibiotique (100 F/ml) → quantité 5",
     "Vérifier le total affiché, enregistrer en statut « Payé »"],
    ["Sous-total traitement = 500 F, Total = 3 500 F",
     "Suggestion de dose affichée si poids + mg/kg renseignés",
     "Stock de l'antibiotique décrémenté de 5",
     "Une vente « 🩺 Consultation » apparaît en Caisse/Historique avec le détail"], badge_color=VERT)

# ══ TEST 5 — Anti double-facturation ══
slide_test(5, "Synchronisation consultation ↔ caisse",
    "Vérifier qu'on ne peut pas facturer deux fois",
    ["Créer une consultation avec traitements en statut « En attente »",
     "Ouvrir Créances → la retrouver au nom du propriétaire → « Marquer payé »",
     "Retourner dans Consultations → vérifier le statut",
     "En Caisse → historique → tenter de SUPPRIMER la vente liée (badge 🩺)"],
    ["La consultation est passée automatiquement à « Payé »",
     "Le stock n'a été décompté qu'UNE seule fois",
     "La suppression est bloquée avec un message explicite",
     "Badge « 🩺 Consultation » visible sur la vente en Caisse"], badge_color=ROUGE)

# ══ TEST 6 — Vaccination ══
slide_test(6, "Vaccination — carnet, rappel, agenda",
    "Valider le cycle vaccinal complet",
    ["Consultation pour un patient existant → traitement = vaccin (catégorie Vaccin)",
     "Le champ « 💉 Rappel vaccinal » apparaît → mettre une date DANS 10 JOURS (pour le test)",
     "Enregistrer en « Payé »",
     "Vérifier : fiche Patient, Dashboard, Agenda"],
    ["Fiche patient : badge « 💉 Vaccin J-10 » + vaccination listée",
     "Dashboard : panneau « Rappels vaccinaux » avec le patient et le téléphone du propriétaire",
     "Agenda : RDV « Vaccination » créé à la date de rappel",
     "Stock de doses décrémenté"], badge_color=RGBColor(0x93, 0x33, 0xEA))

# ══ TEST 7 — Rapports & marges ══
slide_test(7, "Rapports — marges et CA par jour",
    "Valider la rentabilité par produit et l'analyse hebdomadaire",
    ["Ouvrir Rapports (menu Financier) → période « Semaine »",
     "Vérifier le KPI « 📈 Marge brute » et son % du CA",
     "Regarder le Top produits : chaque produit affiche sa marge en F et %",
     "Vérifier le graphique « 📅 CA par jour de la semaine »"],
    ["Marge brute = prix de vente − prix d'achat des ventes payées",
     "Les ventes du test 2 apparaissent avec une marge exacte (pa figé)",
     "Le meilleur jour est surligné en vert",
     "Le bilan affiche : Recettes / Marge brute / Dépenses / Bénéfice net"])

# ══ TEST 8 — Fournisseurs ══
slide_test(8, "Dettes fournisseurs — échéances",
    "Valider l'alerte avant retard de paiement",
    ["Commandes → Nouvelle commande → remplir + « Échéance paiement » = DANS 5 JOURS",
     "Enregistrer, puis marquer la commande « Reçu »",
     "Ouvrir le Dashboard",
     "Fournisseurs → onglet Dettes & Paiements → enregistrer un versement du montant total"],
    ["Dashboard : panneau « 🏭 Échéances fournisseurs » avec badge J-5 et solde dû",
     "Le bouton « Gérer les dettes → » ouvre la bonne page",
     "Après versement complet : l'alerte disparaît du Dashboard"], badge_color=ORANGE)

# ══ TEST 9 — Création de comptes ══
slide_test(9, "Création de comptes (sans rate limit)",
    "Valider l'endpoint admin /api/create-user",
    ["Se connecter en ADMIN → Comptes utilisateurs",
     "Créer 3 comptes d'affilée (ex : test1@labarakat.tg, test2…, test3…)",
     "Se connecter avec l'un d'eux pour vérifier le rôle",
     "Supprimer les comptes de test ensuite"],
    ["Les 3 créations passent SANS erreur « Trop de tentatives »",
     "Chaque compte a le bon rôle dans profiles",
     "Le nouveau compte accède aux pages de son rôle"])

# ══ SLIDE Final — Checklist ══
slide_liste("📋 Checklist de validation finale", [
    "☐  Test 1 — Rôles & sécurité RLS",
    "☐  Test 2 — Caisse : TTC + stock par conditionnement",
    "☐  Test 3 — Créances : versements partiels",
    "☐  Test 4 — Consultation : traitements facturés (5 ml = 500 F)",
    "☐  Test 5 — Anti double-facturation",
    "☐  Test 6 — Vaccination : carnet + rappel + RDV",
    "☐  Test 7 — Rapports : marges + CA par jour",
    "☐  Test 8 — Échéances fournisseurs",
    "☐  Test 9 — Création de comptes",
    "",
    "✅ Tout est vert → basculer l'équipe sur la nouvelle URL et supprimer l'ancien projet Cloudflare",
    "⚠️ Un test échoue → capture d'écran + erreur console (F12) et on corrige ensemble",
], couleur=VERT, taille=17)

prs.save("PLAN_TESTS_LaBarakat.pptx")
print("OK — PLAN_TESTS_LaBarakat.pptx généré,", len(prs.slides.__iter__.__self__._sldIdLst), "diapositives")
